from Markdown2docx import Markdown2docx
from openpyxl import Workbook
from pptx import Presentation

def write_presentation(spreadsheet):
    presentation = Presentation()
    slide_layout = presentation.slide_layouts[1]    # title and content layout
    
    spreadheet = spreadsheet[2:]

    # Build Slides
    for row in spreadheet:
        if(len(row) < 5):
            continue
        slide=presentation.slides.add_slide(slide_layout)

        slide.placeholders[0].text = row[1]
        slide.placeholders[1].text = "Duration " + row[2] + "\nDependencies: " + row[3] + "\nStatus: " + row[4] + "\nResources: " + row[5]
    presentation.save('output/plan_presentation.pptx')


def write_plan(plan):
    # Format as an array
    rows = plan.split('\n')
    spreadsheet = []
    for row in rows:
        split_row = row.split('|')
        spreadsheet.append(split_row)
    
    # Write to Excel
    wb = Workbook()
    ws = wb.active
    for row in spreadsheet:
        ws.append(row)
    wb.save('output/project_plan.xlsx')
    return spreadsheet

def write_document(document, doctype, write_html=True):
    # Save Charter
    if( doctype == 1):
        my_file = open('output/project_charter.md', 'w')
        my_file.write(document)
        my_file.close()
        doc_file = Markdown2docx('output/project_charter')
    # Save Risk Management Plan
    elif(doctype == 2):
        my_file = open('output/risk_management.md', 'w')
        my_file.write(document)
        my_file.close()
        doc_file = Markdown2docx('output/risk_management')

    doc_file.eat_soup()
    if(write_html == True):
        doc_file.write_html()
    doc_file.save()

